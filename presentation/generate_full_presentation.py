from pptx import Presentation
from pptx.util import Inches

def generate_presentation(filename="Pipeline_AI_DSS_full.pptx", summary=None, examples=None):
    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Pipeline AI DSS"
    slide.placeholders[1].text = "Интеллектуальная система поддержки решений — обзор"

    # Architecture
    s = prs.slides.add_slide(prs.slide_layouts[1])
    s.shapes.title.text = "Архитектура"
    s.placeholders[1].text = summary or "Backend: FastAPI; Fuzzy engine; Integrity module; GIS; Frontend."

    # Use cases / Examples
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Пример: ранжирование дефектов"
    s2.placeholders[1].text = examples or "Пример данных и выводов."

    # Recommendations
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "Рекомендации"
    s3.placeholders[1].text = "Дальнейшие шаги для внедрения"

    prs.save(filename)
    return filename
