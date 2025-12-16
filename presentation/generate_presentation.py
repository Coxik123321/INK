from pptx import Presentation

def generate():
    prs = Presentation()

    slides = [
        "Pipeline AI DSS",
        "Актуальность и проблемы",
        "Архитектура системы",
        "Методики и НТД",
        "Экономический эффект",
        "Результаты внедрения",
        "Выводы"
    ]

    for title in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title

    prs.save("Pipeline_AI_DSS.pptx")
